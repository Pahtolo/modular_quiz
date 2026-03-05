from __future__ import annotations

from pathlib import Path
from typing import Iterable

from .ocr import ocr_pdf
from .types import ExtractedMaterial, SourceFile

SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".pptx"}


def _resolve_input_path(input_path: Path | str) -> Path:
    return Path(input_path).expanduser().resolve()


def collect_sources(paths: Iterable[Path | str]) -> tuple[list[SourceFile], list[str]]:
    sources: list[SourceFile] = []
    warnings: list[str] = []
    seen_paths: set[Path] = set()
    seen_warning_messages: set[str] = set()

    def add_source(file_path: Path, source_kind: str) -> None:
        resolved = file_path.resolve()
        if resolved in seen_paths:
            return
        seen_paths.add(resolved)
        sources.append(SourceFile(path=resolved, source_kind=source_kind))

    def warn_once(message: str) -> None:
        if message in seen_warning_messages:
            return
        seen_warning_messages.add(message)
        warnings.append(message)

    for input_path in paths:
        path = _resolve_input_path(input_path)
        if not path.exists():
            warn_once(f"Missing path skipped: {path}")
            continue

        if path.is_file():
            if path.suffix.lower() in SUPPORTED_EXTENSIONS:
                add_source(path, "file")
            else:
                warn_once(f"Unsupported file extension skipped: {path}")
            continue

        if path.is_dir():
            for file_path in sorted(p for p in path.rglob("*") if p.is_file()):
                if file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
                    add_source(file_path, "folder_scan")
                else:
                    warn_once(f"Unsupported file extension skipped: {file_path}")
            continue

        warn_once(f"Unsupported path type skipped: {path}")

    return sources, warnings


def collect_supported_files(paths: Iterable[Path | str]) -> tuple[list[Path], list[str]]:
    sources, warnings = collect_sources(paths)
    return [source.path for source in sources], warnings


def _extract_txt_or_md(path: Path) -> ExtractedMaterial:
    try:
        content = path.read_text(encoding="utf-8", errors="ignore")
        return ExtractedMaterial(path=path, content=content, extracted_by="text")
    except Exception as exc:
        return ExtractedMaterial(path=path, extracted_by="text", errors=[f"Read failed: {exc}"])


def _extract_docx(path: Path) -> ExtractedMaterial:
    try:
        import docx  # python-docx package
    except Exception:
        return ExtractedMaterial(
            path=path,
            extracted_by="docx",
            errors=["python-docx is not installed. Run dependency installer."],
        )

    try:
        doc = docx.Document(str(path))
        parts = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
        content = "\n".join(parts)
        return ExtractedMaterial(path=path, content=content, extracted_by="docx")
    except Exception as exc:
        return ExtractedMaterial(path=path, extracted_by="docx", errors=[f"DOCX extraction failed: {exc}"])


def _extract_pptx(path: Path) -> ExtractedMaterial:
    try:
        from pptx import Presentation
    except Exception:
        return ExtractedMaterial(
            path=path,
            extracted_by="pptx",
            errors=["python-pptx is not installed. Run dependency installer."],
        )

    def _shape_text(shape) -> list[str]:
        parts: list[str] = []

        if hasattr(shape, "has_table") and shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    text = (cell.text or "").strip()
                    if text:
                        parts.append(text)

        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            text_frame = shape.text_frame
            frame_lines = [
                run.text.strip()
                for paragraph in text_frame.paragraphs
                for run in paragraph.runs
                if run.text and run.text.strip()
            ]
            if frame_lines:
                parts.extend(frame_lines)
            else:
                text = (text_frame.text or "").strip()
                if text:
                    parts.append(text)

        fallback = getattr(shape, "text", "")
        if fallback and fallback.strip():
            parts.append(fallback.strip())

        if hasattr(shape, "shapes"):
            for nested in shape.shapes:
                parts.extend(_shape_text(nested))

        return parts

    try:
        prs = Presentation(str(path))
        lines: list[str] = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_lines: list[str] = []
            for shape in slide.shapes:
                slide_lines.extend(_shape_text(shape))

            try:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
            except Exception:
                notes_text = ""
            if notes_text:
                slide_lines.append(f"Notes: {notes_text}")

            deduped: list[str] = []
            seen: set[str] = set()
            for line in slide_lines:
                key = line.strip()
                if not key or key in seen:
                    continue
                seen.add(key)
                deduped.append(key)

            if deduped:
                lines.append(f"Slide {slide_idx}:")
                lines.extend(deduped)
        content = "\n".join(lines).strip()
        return ExtractedMaterial(path=path, content=content, extracted_by="pptx")
    except Exception as exc:
        return ExtractedMaterial(path=path, extracted_by="pptx", errors=[f"PPTX extraction failed: {exc}"])


def _extract_pdf(path: Path, min_chars_for_text: int = 120) -> ExtractedMaterial:
    try:
        from pypdf import PdfReader
    except Exception:
        return ExtractedMaterial(
            path=path,
            extracted_by="pdf",
            needs_ocr=True,
            warnings=["pypdf not installed; marking PDF for OCR fallback."],
        )

    try:
        reader = PdfReader(str(path))
        chunks: list[str] = []
        for page in reader.pages:
            chunks.append((page.extract_text() or "").strip())
        content = "\n\n".join(c for c in chunks if c)
        extracted_text = content.strip()
        needs_ocr = not extracted_text or len(extracted_text) < min_chars_for_text
        warnings: list[str] = []
        if needs_ocr:
            warnings.append("Low PDF text extraction; OCR fallback recommended.")
            try:
                ocr_text, ocr_warnings = ocr_pdf(path)
                warnings.extend(ocr_warnings)
                normalized_ocr = (ocr_text or "").strip()
                if normalized_ocr:
                    merged_text = "\n\n".join(part for part in [extracted_text, normalized_ocr] if part).strip()
                    warnings.append("OCR fallback applied to PDF content.")
                    return ExtractedMaterial(
                        path=path,
                        content=merged_text,
                        extracted_by="pdf+ocr",
                        needs_ocr=False,
                        warnings=warnings,
                    )
                warnings.append("OCR fallback completed but produced no additional text.")
            except Exception as exc:
                warnings.append(f"OCR fallback failed: {exc}")
        return ExtractedMaterial(
            path=path,
            content=content,
            extracted_by="pdf",
            needs_ocr=needs_ocr,
            warnings=warnings,
        )
    except Exception as exc:
        return ExtractedMaterial(
            path=path,
            extracted_by="pdf",
            needs_ocr=True,
            errors=[f"PDF extraction failed: {exc}"],
        )


def extract_material(source: SourceFile) -> ExtractedMaterial:
    ext = source.path.suffix.lower()
    if ext in {".txt", ".md"}:
        return _extract_txt_or_md(source.path)
    if ext == ".docx":
        return _extract_docx(source.path)
    if ext == ".pptx":
        return _extract_pptx(source.path)
    if ext == ".pdf":
        return _extract_pdf(source.path)

    return ExtractedMaterial(
        path=source.path,
        extracted_by="unsupported",
        warnings=[f"Unsupported extension skipped: {source.path}"],
    )


def extract_all_materials(sources: list[SourceFile]) -> list[ExtractedMaterial]:
    return [extract_material(src) for src in sources]


def extract_paths(paths: Iterable[Path | str]) -> tuple[list[ExtractedMaterial], list[str]]:
    sources, collection_warnings = collect_sources(paths)
    return extract_all_materials(sources), collection_warnings
