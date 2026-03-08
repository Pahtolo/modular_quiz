from __future__ import annotations

import json
import tempfile
import unittest
from pathlib import Path
from unittest.mock import MagicMock, patch

from fastapi.testclient import TestClient

from quiz_app.api.server import create_app
from quiz_app.models import MCQQuestion, Quiz, ShortQuestion
from quiz_app.openai_auth import OAuthTokenSet
from quiz_app.providers import ModelOption


class _StubProvider:
    provider_name = "stub"

    def generate_quiz(
        self,
        materials_text: str,
        title_hint: str,
        instructions_hint: str,
        total_questions: int,
        mcq_count: int,
        short_count: int,
        mcq_options: int,
        model: str | None = None,
    ) -> Quiz:
        self._assertions(
            materials_text,
            total_questions,
            mcq_count,
            short_count,
            mcq_options,
        )
        return Quiz(
            title=title_hint or "Generated Quiz",
            instructions=instructions_hint or "Answer all questions.",
            questions=[
                MCQQuestion(
                    id="q1",
                    prompt="Pick A",
                    points=1,
                    options=["A", "B", "C", "D"],
                    answer="A",
                ),
                ShortQuestion(
                    id="q2",
                    prompt="Say hi",
                    points=2,
                    expected="hi",
                ),
            ],
        )

    @staticmethod
    def _assertions(
        materials_text: str,
        total_questions: int,
        mcq_count: int,
        short_count: int,
        mcq_options: int,
    ) -> None:
        assert "Generation template (JSON):" in materials_text
        assert '"questions"' in materials_text
        assert "Source material:" in materials_text
        assert "hello from txt" in materials_text
        assert total_questions == 2
        assert mcq_count == 1
        assert short_count == 1
        assert mcq_options == 4


class _ModelAwareStubProvider(_StubProvider):
    def __init__(self, available_models: list[str]):
        self.available_models = available_models
        self.last_generation_model: str | None = None

    def list_models(self) -> list[ModelOption]:
        return [
            ModelOption(
                id=model_id,
                label=model_id,
                provider="claude",
                capability_tags=("generation",),
            )
            for model_id in self.available_models
        ]

    def generate_quiz(
        self,
        materials_text: str,
        title_hint: str,
        instructions_hint: str,
        total_questions: int,
        mcq_count: int,
        short_count: int,
        mcq_options: int,
        model: str | None = None,
    ) -> Quiz:
        self.last_generation_model = model
        return super().generate_quiz(
            materials_text=materials_text,
            title_hint=title_hint,
            instructions_hint=instructions_hint,
            total_questions=total_questions,
            mcq_count=mcq_count,
            short_count=short_count,
            mcq_options=mcq_options,
            model=model,
        )


class APIServerTests(unittest.TestCase):
    def setUp(self) -> None:
        self.tmp = tempfile.TemporaryDirectory()
        self.root = Path(self.tmp.name)

        settings_dir = self.root / "settings"
        settings_dir.mkdir(parents=True, exist_ok=True)
        self.settings_path = settings_dir / "settings.json"

        self.api_token = "test-token"
        app = create_app(
            settings_path=self.settings_path,
            api_token=self.api_token,
            project_root=self.root,
        )
        self.client = TestClient(app)
        self.headers = {"Authorization": f"Bearer {self.api_token}"}

    def tearDown(self) -> None:
        self.tmp.cleanup()

    def _post(self, path: str, body: dict) -> dict:
        response = self.client.post(path, headers=self.headers, json=body)
        self.assertEqual(response.status_code, 200, response.text)
        return response.json()

    def test_health_requires_auth(self) -> None:
        unauthorized = self.client.get("/v1/health")
        self.assertEqual(unauthorized.status_code, 401)

        ok = self.client.get("/v1/health", headers=self.headers)
        self.assertEqual(ok.status_code, 200)
        self.assertTrue(ok.json()["ok"])

    def test_settings_round_trip(self) -> None:
        initial = self.client.get("/v1/settings", headers=self.headers)
        self.assertEqual(initial.status_code, 200)
        settings = initial.json()["settings"]
        self.assertIn("quiz_roots", settings)
        self.assertEqual(settings["quiz_clock_mode"], "stopwatch")
        self.assertEqual(settings["quiz_timer_duration_seconds"], 900)
        self.assertFalse(settings["shuffle_mcq_answers"])

        updated = self.client.put(
            "/v1/settings",
            headers=self.headers,
            json={
                "quiz_dir": str(self.root / "Algorithm Analysis"),
                "quiz_roots": [str(self.root)],
                "feedback_mode": "end_only",
                "quiz_clock_mode": "timer",
                "quiz_timer_duration_seconds": 1200,
                "question_timer_seconds": 45,
                "lock_questions_by_progression": False,
                "shuffle_mcq_answers": True,
                "generation_defaults": {
                    "total": 10,
                    "mcq_count": 7,
                    "short_count": 3,
                    "mcq_options": 4,
                },
            },
        )
        self.assertEqual(updated.status_code, 200)
        payload = updated.json()["settings"]
        self.assertEqual(payload["feedback_mode"], "end_only")
        self.assertFalse(payload["show_feedback_on_answer"])
        self.assertTrue(payload["show_feedback_on_completion"])
        self.assertFalse(payload["auto_advance_enabled"])
        self.assertEqual(payload["quiz_clock_mode"], "timer")
        self.assertEqual(payload["quiz_timer_duration_seconds"], 1200)
        self.assertEqual(payload["question_timer_seconds"], 45)
        self.assertFalse(payload["lock_questions_by_progression"])
        self.assertTrue(payload["shuffle_mcq_answers"])
        self.assertEqual(payload["quiz_roots"], [str((self.root / "Quizzes").resolve())])

        clock_off = self.client.put(
            "/v1/settings",
            headers=self.headers,
            json={
                "quiz_clock_mode": "off",
            },
        )
        self.assertEqual(clock_off.status_code, 200)
        clock_off_payload = clock_off.json()["settings"]
        self.assertEqual(clock_off_payload["quiz_clock_mode"], "off")
        self.assertEqual(clock_off_payload["quiz_timer_duration_seconds"], 1200)

        flags_update = self.client.put(
            "/v1/settings",
            headers=self.headers,
            json={
                "show_feedback_on_answer": True,
                "show_feedback_on_completion": False,
                "auto_advance_enabled": True,
            },
        )
        self.assertEqual(flags_update.status_code, 200)
        flags_payload = flags_update.json()["settings"]
        self.assertEqual(flags_payload["feedback_mode"], "auto_advance")
        self.assertTrue(flags_payload["show_feedback_on_answer"])
        self.assertFalse(flags_payload["show_feedback_on_completion"])
        self.assertTrue(flags_payload["auto_advance_enabled"])

    def test_settings_sanitizes_deprecated_claude_models(self) -> None:
        updated = self.client.put(
            "/v1/settings",
            headers=self.headers,
            json={
                "claude_model_selected": "claude-3-7-sonnet-latest",
                "claude_models": ["claude-3-7-sonnet-latest", "claude-3-5-haiku-latest"],
                "preferred_model_key": "claude:claude-3-7-sonnet-latest",
            },
        )
        self.assertEqual(updated.status_code, 200, updated.text)
        payload = updated.json()["settings"]
        self.assertNotIn("claude-3-7-sonnet-latest", payload["claude_models"])
        self.assertNotEqual(payload["claude_model_selected"], "claude-3-7-sonnet-latest")
        self.assertNotEqual(payload["preferred_model_key"], "claude:claude-3-7-sonnet-latest")

    def test_models_preview_lists_models_from_draft_settings(self) -> None:
        preview_provider = MagicMock()
        preview_provider.list_models.return_value = [
            ModelOption(
                id="claude-preview",
                label="Claude Preview",
                provider="claude",
                capability_tags=("generation",),
            )
        ]
        with patch("quiz_app.api.server._provider_client_for_preview", return_value=preview_provider):
            response = self.client.post(
                "/v1/models/preview",
                headers=self.headers,
                json={
                    "provider": "claude",
                    "settings": {
                        "claude_api_key": "draft-key",
                    },
                },
            )
        self.assertEqual(response.status_code, 200, response.text)
        payload = response.json()
        self.assertEqual(payload["models"][0]["id"], "claude-preview")
        preview_provider.list_models.assert_called_once()

    def test_quiz_tree_and_load(self) -> None:
        quiz_dir = self.root / "Algorithm Analysis"
        quiz_dir.mkdir(parents=True, exist_ok=True)
        quiz_path = quiz_dir / "sample.json"
        nested_dir = quiz_dir / "week1"
        nested_dir.mkdir(parents=True, exist_ok=True)
        nested_quiz_path = nested_dir / "nested.json"
        ignored_file = nested_dir / "notes.txt"
        ignored_file.write_text("ignore me", encoding="utf-8")

        quiz_path.write_text(
            json.dumps(
                {
                    "title": "Sample Quiz",
                    "instructions": "Answer all questions.",
                    "questions": [
                        {
                            "type": "mcq",
                            "id": "q1",
                            "prompt": "2+2?",
                            "options": ["3", "4"],
                            "answer": "B",
                            "points": 1,
                        }
                    ],
                }
            ),
            encoding="utf-8",
        )
        nested_quiz_path.write_text(
            json.dumps(
                {
                    "title": "Nested Quiz",
                    "instructions": "Answer all questions.",
                    "questions": [
                        {
                            "type": "mcq",
                            "id": "q1",
                            "prompt": "1+1?",
                            "options": ["1", "2"],
                            "answer": "B",
                            "points": 1,
                        }
                    ],
                }
            ),
            encoding="utf-8",
        )

        self.client.put(
            "/v1/settings",
            headers=self.headers,
            json={"quiz_roots": [str(self.root)]},
        )

        tree = self._post("/v1/quizzes/tree", {"quiz_roots": [str(self.root)]})
        self.assertTrue(tree["roots"])
        root_node = tree["roots"][0]
        self.assertEqual(root_node["kind"], "root")

        def _collect(nodes: list[dict], kind: str, field: str) -> set[str]:
            values: set[str] = set()
            for node in nodes:
                if node.get("kind") == kind:
                    value = node.get(field)
                    if isinstance(value, str):
                        values.add(value)
                values.update(_collect(node.get("children", []), kind, field))
            return values

        folder_names = _collect(root_node["children"], "folder", "name")
        self.assertIn("Algorithm Analysis", folder_names)
        self.assertIn("week1", folder_names)

        child_names = _collect(root_node["children"], "quiz", "name")
        self.assertIn("Sample Quiz", child_names)
        self.assertIn("Nested Quiz", child_names)

        child_rel_paths = _collect(root_node["children"], "quiz", "relative_path")
        self.assertIn("Algorithm Analysis/sample.json", child_rel_paths)
        self.assertIn("Algorithm Analysis/week1/nested.json", child_rel_paths)
        self.assertNotIn("Algorithm Analysis/week1/notes.txt", child_rel_paths)

        loaded = self._post("/v1/quizzes/load", {"path": str(quiz_path)})
        self.assertEqual(loaded["quiz"]["title"], "Sample Quiz")

        invalid = self.client.post(
            "/v1/quizzes/load",
            headers=self.headers,
            json={"path": str(self.root / "missing.json")},
        )
        self.assertEqual(invalid.status_code, 404)

    @patch("quiz_app.api.server.OpenAIPKCEAuthenticator.authorize_in_browser")
    def test_openai_oauth_connect_persists_tokens(self, mock_authorize) -> None:
        mock_authorize.return_value = OAuthTokenSet(
            access_token="access-token-123",
            refresh_token="refresh-token-123",
            expires_at=1234567890.0,
            token_type="Bearer",
        )
        update = self.client.put(
            "/v1/settings",
            headers=self.headers,
            json={
                "openai_oauth_client_id": "client-id-123",
            },
        )
        self.assertEqual(update.status_code, 200, update.text)

        response = self.client.post(
            "/v1/oauth/openai/connect",
            headers=self.headers,
            json={},
        )
        self.assertEqual(response.status_code, 200, response.text)
        payload = response.json()
        self.assertTrue(payload["ok"])
        self.assertEqual(payload["token_type"], "Bearer")
        self.assertEqual(payload["expires_at"], 1234567890.0)
        self.assertTrue(payload["refresh_token_present"])

        settings_response = self.client.get("/v1/settings", headers=self.headers)
        self.assertEqual(settings_response.status_code, 200, settings_response.text)
        settings = settings_response.json()["settings"]
        self.assertEqual(settings["openai_auth_mode"], "oauth")
        self.assertEqual(settings["openai_oauth_access_token"], "access-token-123")
        self.assertEqual(settings["openai_oauth_refresh_token"], "refresh-token-123")
        self.assertEqual(settings["openai_oauth_expires_at"], 1234567890.0)

    def test_openai_oauth_connect_requires_client_id(self) -> None:
        response = self.client.post(
            "/v1/oauth/openai/connect",
            headers=self.headers,
            json={},
        )
        self.assertEqual(response.status_code, 422, response.text)
        payload = response.json()["error"]
        self.assertEqual(payload["code"], "VALIDATION_ERROR")
        self.assertIn("client ID", payload["message"])

    def test_quiz_tree_multiple_roots_distinct_paths(self) -> None:
        root_one = self.root / "coursesA" / "quizzes"
        root_two = self.root / "coursesB" / "quizzes"
        root_one.mkdir(parents=True, exist_ok=True)
        root_two.mkdir(parents=True, exist_ok=True)

        (root_one / "a.json").write_text(
            json.dumps(
                {
                    "title": "A",
                    "instructions": "",
                    "questions": [
                        {
                            "type": "mcq",
                            "id": "q1",
                            "prompt": "A?",
                            "options": ["A", "B"],
                            "answer": "A",
                        }
                    ],
                }
            ),
            encoding="utf-8",
        )
        (root_two / "b.json").write_text(
            json.dumps(
                {
                    "title": "B",
                    "instructions": "",
                    "questions": [
                        {
                            "type": "mcq",
                            "id": "q1",
                            "prompt": "B?",
                            "options": ["A", "B"],
                            "answer": "B",
                        }
                    ],
                }
            ),
            encoding="utf-8",
        )

        tree = self._post(
            "/v1/quizzes/tree",
            {"quiz_roots": [str(root_one), str(root_two)]},
        )
        self.assertEqual(len(tree["roots"]), 2)
        root_names = {node["name"] for node in tree["roots"]}
        self.assertEqual(root_names, {str(root_one.resolve()), str(root_two.resolve())})
        for node in tree["roots"]:
            self.assertEqual(node["kind"], "root")
            self.assertEqual(len(node["children"]), 1)
            self.assertEqual(node["children"][0]["kind"], "quiz")
            self.assertTrue(node["children"][0]["path"].lower().endswith(".json"))

    def test_quizzes_library_import_and_structure(self) -> None:
        source_folder = self.root / "source-quizzes"
        nested = source_folder / "week1"
        nested.mkdir(parents=True, exist_ok=True)

        (source_folder / "intro.json").write_text(
            json.dumps(
                {
                    "name": "Intro Display Name",
                    "title": "Intro",
                    "instructions": "",
                    "questions": [
                        {
                            "type": "mcq",
                            "id": "q1",
                            "prompt": "A?",
                            "options": ["A", "B"],
                            "answer": "A",
                        }
                    ],
                }
            ),
            encoding="utf-8",
        )
        (nested / "deep.json").write_text(
            json.dumps(
                {
                    "title": "Deep",
                    "instructions": "",
                    "questions": [
                        {
                            "type": "mcq",
                            "id": "q1",
                            "prompt": "B?",
                            "options": ["A", "B"],
                            "answer": "B",
                        }
                    ],
                }
            ),
            encoding="utf-8",
        )
        (nested / "ignore.txt").write_text("ignore", encoding="utf-8")

        library = self.client.get("/v1/quizzes/library", headers=self.headers)
        self.assertEqual(library.status_code, 200, library.text)
        quizzes_dir = Path(library.json()["quizzes_dir"])
        self.assertTrue(quizzes_dir.exists())
        self.assertEqual(quizzes_dir, (self.root / "Quizzes").resolve())
        self.assertEqual(library.json()["settings"]["quiz_roots"], [str(quizzes_dir)])

        imported = self._post(
            "/v1/quizzes/library/import",
            {"source_paths": [str(source_folder)]},
        )
        self.assertEqual(imported["imported_files"], 2)
        self.assertFalse(imported["warnings"], imported["warnings"])
        self.assertEqual(imported["settings"]["quiz_roots"], [str(quizzes_dir)])

        def _collect_quiz_names(nodes: list[dict]) -> list[str]:
            names: list[str] = []
            for node in nodes:
                if node.get("kind") == "quiz":
                    names.append(str(node.get("name", "")))
                names.extend(_collect_quiz_names(node.get("children", [])))
            return names

        quiz_names = _collect_quiz_names(imported["tree"])
        self.assertIn("Intro", quiz_names)
        self.assertIn("Deep", quiz_names)

        tree_default = self._post("/v1/quizzes/tree", {})
        self.assertTrue(tree_default["roots"])
        default_root = tree_default["roots"][0]
        self.assertEqual(default_root["path"], str(quizzes_dir))

        def _collect_quiz_rel_paths(nodes: list[dict]) -> set[str]:
            values: set[str] = set()
            for node in nodes:
                if node.get("kind") == "quiz":
                    rel = node.get("relative_path")
                    if isinstance(rel, str):
                        values.add(rel)
                values.update(_collect_quiz_rel_paths(node.get("children", [])))
            return values

        default_paths = _collect_quiz_rel_paths(default_root["children"])
        self.assertIn(f"{source_folder.name}/intro.json", default_paths)
        self.assertIn(f"{source_folder.name}/week1/deep.json", default_paths)

    def test_quizzes_library_uses_settings_root_not_project_root(self) -> None:
        settings_root = self.root / "user-data"
        settings_path = settings_root / "settings" / "settings.json"
        settings_path.parent.mkdir(parents=True, exist_ok=True)

        project_root = self.root / "packaged-app-root"
        bundled_quizzes = project_root / "Quizzes"
        bundled_quizzes.mkdir(parents=True, exist_ok=True)
        (bundled_quizzes / "starter.json").write_text(
            json.dumps(
                {
                    "title": "Bundled Starter Quiz",
                    "instructions": "",
                    "questions": [
                        {
                            "type": "mcq",
                            "id": "q1",
                            "prompt": "A?",
                            "options": ["A", "B"],
                            "answer": "A",
                        }
                    ],
                }
            ),
            encoding="utf-8",
        )

        token = "isolated-token"
        app = create_app(
            settings_path=settings_path,
            api_token=token,
            project_root=project_root,
        )
        client = TestClient(app)
        headers = {"Authorization": f"Bearer {token}"}

        response = client.get("/v1/quizzes/library", headers=headers)
        self.assertEqual(response.status_code, 200, response.text)
        payload = response.json()

        expected_dir = (settings_root / "Quizzes").resolve()
        self.assertEqual(Path(payload["quizzes_dir"]), expected_dir)
        self.assertEqual(payload["settings"]["quiz_roots"], [str(expected_dir)])
        self.assertEqual(payload["tree"], [])
        self.assertFalse((expected_dir / "starter.json").exists())

    def test_quizzes_library_rename_updates_json_title(self) -> None:
        quizzes_dir = self.root / "Quizzes"
        quizzes_dir.mkdir(parents=True, exist_ok=True)
        quiz_path = quizzes_dir / "rename-me.json"
        quiz_path.write_text(
            json.dumps(
                {
                    "title": "Old Title",
                    "instructions": "Answer all questions.",
                    "questions": [
                        {
                            "type": "mcq",
                            "id": "q1",
                            "prompt": "A?",
                            "options": ["A", "B"],
                            "answer": "A",
                        }
                    ],
                }
            ),
            encoding="utf-8",
        )

        renamed = self._post(
            "/v1/quizzes/library/rename",
            {
                "path": str(quiz_path),
                "title": "New Title",
            },
        )
        self.assertEqual(renamed["path"], str(quiz_path.resolve()))
        self.assertEqual(renamed["title"], "New Title")

        saved = json.loads(quiz_path.read_text(encoding="utf-8"))
        self.assertEqual(saved["title"], "New Title")

        tree = self._post("/v1/quizzes/tree", {})
        root = tree["roots"][0]
        quiz_titles = [node["name"] for node in root["children"] if node.get("kind") == "quiz"]
        self.assertIn("New Title", quiz_titles)

    def test_quizzes_library_can_create_and_delete_folders_and_files(self) -> None:
        library = self.client.get("/v1/quizzes/library", headers=self.headers)
        self.assertEqual(library.status_code, 200, library.text)
        quizzes_dir = Path(library.json()["quizzes_dir"])

        created = self._post(
            "/v1/quizzes/library/folder",
            {
                "name": "Week 4",
            },
        )
        week_folder = quizzes_dir / "Week 4"
        self.assertEqual(created["path"], str(week_folder.resolve()))
        self.assertTrue(week_folder.is_dir())
        self.assertTrue(any(node.get("path") == str(week_folder.resolve()) for node in created["tree"]))

        nested = self._post(
            "/v1/quizzes/library/folder",
            {
                "name": "Practice",
                "parent_path": str(week_folder),
            },
        )
        practice_folder = week_folder / "Practice"
        self.assertEqual(nested["path"], str(practice_folder.resolve()))
        self.assertTrue(practice_folder.is_dir())

        quiz_path = practice_folder / "sample.json"
        quiz_path.write_text(
            json.dumps(
                {
                    "title": "Sample Quiz",
                    "instructions": "",
                    "questions": [
                        {
                            "type": "mcq",
                            "id": "q1",
                            "prompt": "A?",
                            "options": ["A", "B"],
                            "answer": "A",
                        }
                    ],
                }
            ),
            encoding="utf-8",
        )
        self.assertTrue(quiz_path.exists())

        deleted_quiz = self._post(
            "/v1/quizzes/library/delete",
            {
                "path": str(quiz_path),
            },
        )
        self.assertEqual(deleted_quiz["deleted_kind"], "quiz")
        self.assertFalse(quiz_path.exists())

        (practice_folder / "sample.json").write_text(
            json.dumps(
                {
                    "title": "Sample Quiz",
                    "instructions": "",
                    "questions": [
                        {
                            "type": "mcq",
                            "id": "q1",
                            "prompt": "A?",
                            "options": ["A", "B"],
                            "answer": "A",
                        }
                    ],
                }
            ),
            encoding="utf-8",
        )
        deleted_folder = self._post(
            "/v1/quizzes/library/delete",
            {
                "path": str(week_folder),
            },
        )
        self.assertEqual(deleted_folder["deleted_kind"], "folder")
        self.assertFalse(week_folder.exists())
        self.assertEqual(deleted_folder["tree"], [])

    def test_quizzes_library_can_rename_folder_and_move_quiz(self) -> None:
        library = self.client.get("/v1/quizzes/library", headers=self.headers)
        self.assertEqual(library.status_code, 200, library.text)
        quizzes_dir = Path(library.json()["quizzes_dir"])

        source_folder = quizzes_dir / "Week 4"
        source_folder.mkdir(parents=True, exist_ok=True)
        quiz_path = source_folder / "sample.json"
        quiz_path.write_text(
            json.dumps(
                {
                    "title": "Sample Quiz",
                    "instructions": "Answer all questions.",
                    "questions": [
                        {
                            "type": "mcq",
                            "id": "q1",
                            "prompt": "A?",
                            "options": ["A", "B"],
                            "answer": "A",
                        }
                    ],
                }
            ),
            encoding="utf-8",
        )

        renamed = self._post(
            "/v1/quizzes/library/folder/rename",
            {
                "path": str(source_folder),
                "name": "Week Four",
            },
        )
        renamed_folder = quizzes_dir / "Week Four"
        self.assertEqual(renamed["path"], str(renamed_folder.resolve()))
        self.assertTrue(renamed_folder.is_dir())
        self.assertFalse(source_folder.exists())

        destination_folder = quizzes_dir / "Archive"
        destination_folder.mkdir(parents=True, exist_ok=True)
        moved = self._post(
            "/v1/quizzes/library/move",
            {
                "path": str((renamed_folder / "sample.json").resolve()),
                "destination_parent_path": str(destination_folder.resolve()),
            },
        )
        moved_path = destination_folder / "sample.json"
        self.assertEqual(moved["path"], str(moved_path.resolve()))
        self.assertTrue(moved_path.exists())

    def test_quizzes_library_move_rejects_invalid_folder_move(self) -> None:
        library = self.client.get("/v1/quizzes/library", headers=self.headers)
        self.assertEqual(library.status_code, 200, library.text)
        quizzes_dir = Path(library.json()["quizzes_dir"])
        source_folder = quizzes_dir / "Parent"
        child_folder = source_folder / "Child"
        child_folder.mkdir(parents=True, exist_ok=True)

        response = self.client.post(
            "/v1/quizzes/library/move",
            headers=self.headers,
            json={
                "path": str(source_folder.resolve()),
                "destination_parent_path": str(child_folder.resolve()),
            },
        )
        self.assertEqual(response.status_code, 422, response.text)
        self.assertEqual(response.json()["error"]["code"], "VALIDATION_ERROR")

    def test_load_quiz_rejects_duplicate_question_ids(self) -> None:
        quiz_path = self.root / "duplicates.json"
        quiz_path.write_text(
            json.dumps(
                {
                    "title": "Duplicate IDs",
                    "instructions": "Answer all questions.",
                    "questions": [
                        {
                            "type": "mcq",
                            "id": "q1",
                            "prompt": "A?",
                            "options": ["A", "B"],
                            "answer": "A",
                        },
                        {
                            "type": "short",
                            "id": "q1",
                            "prompt": "Explain A",
                            "expected": "A",
                        },
                    ],
                }
            ),
            encoding="utf-8",
        )

        response = self.client.post(
            "/v1/quizzes/load",
            headers=self.headers,
            json={"path": str(quiz_path.resolve())},
        )
        self.assertEqual(response.status_code, 422, response.text)
        self.assertIn("Question IDs must be unique", response.json()["error"]["message"])

    def test_load_quiz_rejects_mcq_with_more_than_four_options(self) -> None:
        quiz_path = self.root / "too-many-options.json"
        quiz_path.write_text(
            json.dumps(
                {
                    "title": "Too Many Options",
                    "instructions": "Answer all questions.",
                    "questions": [
                        {
                            "type": "mcq",
                            "id": "q1",
                            "prompt": "Pick one",
                            "options": ["A", "B", "C", "D", "E"],
                            "answer": "A",
                        }
                    ],
                }
            ),
            encoding="utf-8",
        )

        response = self.client.post(
            "/v1/quizzes/load",
            headers=self.headers,
            json={"path": str(quiz_path.resolve())},
        )
        self.assertEqual(response.status_code, 422, response.text)
        self.assertIn("at most 4 choices", response.json()["error"]["message"])

    def test_grading_endpoints(self) -> None:
        mcq = self._post(
            "/v1/grade/mcq",
            {
                "question": {
                    "id": "q1",
                    "prompt": "2+2?",
                    "options": ["3", "4"],
                    "answer": "B",
                    "points": 1,
                },
                "user_answer": "B",
            },
        )
        self.assertTrue(mcq["result"]["correct"])

        short_self = self._post(
            "/v1/grade/short",
            {
                "provider": "self",
                "question": {
                    "id": "q2",
                    "prompt": "Say hi",
                    "expected": "hi",
                    "points": 2,
                },
                "user_answer": "hello",
                "self_score": 1,
            },
        )
        self.assertEqual(short_self["result"]["points_awarded"], 1)

    def test_grade_mcq_rejects_whitespace_only_options(self) -> None:
        response = self.client.post(
            "/v1/grade/mcq",
            headers=self.headers,
            json={
                "question": {
                    "id": "q1",
                    "prompt": "2+2?",
                    "options": ["   ", "4"],
                    "answer": "B",
                    "points": 1,
                },
                "user_answer": "A",
            },
        )

        self.assertEqual(response.status_code, 422, response.text)
        payload = response.json()["error"]
        self.assertEqual(payload["code"], "VALIDATION_ERROR")
        self.assertIn("non-empty", payload["message"])

    def test_grade_mcq_rejects_non_string_options(self) -> None:
        response = self.client.post(
            "/v1/grade/mcq",
            headers=self.headers,
            json={
                "question": {
                    "id": "q1",
                    "prompt": "2+2?",
                    "options": [1, "4"],
                    "answer": "B",
                    "points": 1,
                },
                "user_answer": "A",
            },
        )

        self.assertEqual(response.status_code, 422, response.text)
        payload = response.json()["error"]
        self.assertEqual(payload["code"], "VALIDATION_ERROR")
        self.assertIn("strings", payload["message"])

    def test_grade_mcq_rejects_out_of_range_answer(self) -> None:
        response = self.client.post(
            "/v1/grade/mcq",
            headers=self.headers,
            json={
                "question": {
                    "id": "q1",
                    "prompt": "2+2?",
                    "options": ["3", "4"],
                    "answer": "Z",
                    "points": 1,
                },
                "user_answer": "A",
            },
        )

        self.assertEqual(response.status_code, 422, response.text)
        payload = response.json()["error"]
        self.assertEqual(payload["code"], "VALIDATION_ERROR")
        self.assertIn("A, B", payload["message"])

    def test_grade_mcq_rejects_multi_character_answer(self) -> None:
        response = self.client.post(
            "/v1/grade/mcq",
            headers=self.headers,
            json={
                "question": {
                    "id": "q1",
                    "prompt": "2+2?",
                    "options": ["3", "4"],
                    "answer": "AB",
                    "points": 1,
                },
                "user_answer": "A",
            },
        )

        self.assertEqual(response.status_code, 422, response.text)
        payload = response.json()["error"]
        self.assertEqual(payload["code"], "VALIDATION_ERROR")
        self.assertIn("exactly one letter", payload["message"])

    def test_grade_mcq_rejects_non_positive_points(self) -> None:
        for points in (0, -3):
            with self.subTest(points=points):
                response = self.client.post(
                    "/v1/grade/mcq",
                    headers=self.headers,
                    json={
                        "question": {
                            "id": "q1",
                            "prompt": "2+2?",
                            "options": ["3", "4"],
                            "answer": "B",
                            "points": points,
                        },
                        "user_answer": "A",
                    },
                )

                self.assertEqual(response.status_code, 422, response.text)
                payload = response.json()["error"]
                self.assertEqual(payload["code"], "VALIDATION_ERROR")
                self.assertIn("positive integer", payload["message"])

    def test_feedback_chat_endpoint(self) -> None:
        provider = MagicMock()
        provider.feedback_chat.return_value = "You should compare your answer to the expected terms."
        with patch("quiz_app.api.server._provider_client", return_value=provider):
            response = self._post(
                "/v1/feedback/chat",
                {
                    "provider": "openai",
                    "model": "gpt-5-mini",
                    "user_message": "Why was my answer marked incorrect?",
                    "feedback": "You are incorrect.",
                    "question": {
                        "id": "q1",
                        "type": "short",
                        "prompt": "Define asymptotic notation.",
                        "options": [],
                    },
                    "user_answer": "It means average runtime.",
                    "expected_answer": "It describes growth rate.",
                    "chat_history": [
                        {"role": "assistant", "text": "You are incorrect."},
                    ],
                },
            )
        self.assertIn("text", response)
        provider.feedback_chat.assert_called_once()

    def test_explain_short_endpoint(self) -> None:
        provider = MagicMock()
        provider.explain_short.return_value = "You are close, but your answer missed the growth-rate detail."
        with patch("quiz_app.api.server._provider_client", return_value=provider):
            response = self._post(
                "/v1/explain/short",
                {
                    "provider": "openai",
                    "model": "gpt-5-mini",
                    "question": {
                        "id": "q1",
                        "type": "short",
                        "prompt": "Define asymptotic notation.",
                        "expected": "It describes algorithm growth rate.",
                        "points": 2,
                    },
                    "user_answer": "It means average runtime.",
                },
            )
        self.assertIn("text", response)
        provider.explain_short.assert_called_once()

    def test_history_append_and_filter(self) -> None:
        record = {
            "timestamp": "2026-03-04T10:00:00",
            "quiz_path": str(self.root / "q.json"),
            "quiz_title": "Q",
            "score": 3,
            "max_score": 4,
            "percent": 75.0,
            "duration_seconds": 12.0,
            "model_key": "self:",
            "quiz_clock_mode": "timer",
            "quiz_timer_duration_seconds": 900,
            "questions": [
                {
                    "question_id": "q1",
                    "question_type": "mcq",
                    "user_answer": "A",
                    "correct_answer_or_expected": "A",
                    "points_awarded": 1,
                    "max_points": 1,
                    "feedback": "Correct.",
                }
            ],
        }
        append_resp = self._post("/v1/history/append", record)
        self.assertTrue(append_resp["ok"])

        history = self.client.get("/v1/history", headers=self.headers)
        self.assertEqual(history.status_code, 200)
        records = history.json()["records"]
        self.assertEqual(len(records), 1)
        self.assertEqual(records[0]["quiz_clock_mode"], "timer")
        self.assertEqual(records[0]["quiz_timer_duration_seconds"], 900)

        filtered = self.client.get(
            "/v1/history",
            headers=self.headers,
            params={"quiz_path": str(self.root / "q.json")},
        )
        self.assertEqual(filtered.status_code, 200)
        self.assertEqual(len(filtered.json()["records"]), 1)

    def test_history_update(self) -> None:
        base_record = {
            "timestamp": "2026-03-04T11:00:00",
            "quiz_path": str(self.root / "q.json"),
            "quiz_title": "Q",
            "score": 0,
            "max_score": 4,
            "percent": 0.0,
            "duration_seconds": 20.0,
            "model_key": "self:",
            "quiz_clock_mode": "stopwatch",
            "quiz_timer_duration_seconds": 0,
            "questions": [
                {
                    "question_id": "q2",
                    "question_type": "short",
                    "user_answer": "hello",
                    "correct_answer_or_expected": "hi",
                    "points_awarded": 0,
                    "max_points": 2,
                    "feedback": "No model selected. Response recorded as ungraded.",
                    "ungraded": True,
                }
            ],
        }
        self._post("/v1/history/append", base_record)

        updated_record = {
            **base_record,
            "score": 2,
            "percent": 50.0,
            "model_key": "openai:gpt-5-mini",
            "quiz_clock_mode": "timer",
            "quiz_timer_duration_seconds": 1200,
            "questions": [
                {
                    "question_id": "q2",
                    "question_type": "short",
                    "user_answer": "hello",
                    "correct_answer_or_expected": "hi",
                    "points_awarded": 2,
                    "max_points": 2,
                    "feedback": "You are correct.",
                    "ungraded": False,
                }
            ],
        }
        updated = self._post(
            "/v1/history/update",
            {
                "match": {
                    "timestamp": base_record["timestamp"],
                    "quiz_path": base_record["quiz_path"],
                    "model_key": base_record["model_key"],
                    "score": base_record["score"],
                    "max_score": base_record["max_score"],
                    "duration_seconds": base_record["duration_seconds"],
                },
                "record": updated_record,
            },
        )
        self.assertTrue(updated["ok"])

        history = self.client.get("/v1/history", headers=self.headers)
        self.assertEqual(history.status_code, 200)
        records = history.json()["records"]
        self.assertEqual(len(records), 1)
        self.assertEqual(records[0]["score"], 2)
        self.assertEqual(records[0]["model_key"], "openai:gpt-5-mini")
        self.assertEqual(records[0]["quiz_clock_mode"], "timer")
        self.assertEqual(records[0]["quiz_timer_duration_seconds"], 1200)
        self.assertEqual(records[0]["questions"][0]["points_awarded"], 2)
        self.assertFalse(records[0]["questions"][0]["ungraded"])

    def test_collect_sources_and_generate_run(self) -> None:
        docs = self.root / "docs"
        docs.mkdir(parents=True, exist_ok=True)
        txt_path = docs / "notes.txt"
        txt_path.write_text("hello from txt", encoding="utf-8")
        unsupported = docs / "skip.bin"
        unsupported.write_bytes(b"\x00\x01")

        pdf_path = docs / "blank.pdf"
        try:
            from pypdf import PdfWriter

            writer = PdfWriter()
            writer.add_blank_page(width=72, height=72)
            with pdf_path.open("wb") as fh:
                writer.write(fh)
        except Exception:
            pdf_path.write_bytes(b"%PDF-1.4\n1 0 obj\n<<>>\nendobj\ntrailer\n<<>>\n%%EOF")

        collect = self._post(
            "/v1/generate/collect-sources",
            {"paths": [str(docs)]},
        )
        self.assertGreaterEqual(len(collect["sources"]), 2)
        self.assertTrue(any("Unsupported file extension" in w for w in collect["warnings"]))

        with patch("quiz_app.api.server._provider_client", return_value=_StubProvider()):
            generated = self._post(
                "/v1/generate/run",
                {
                    "quiz_dir": str(self.root),
                    "sources": collect["sources"],
                    "provider": "claude",
                    "model": "stub-model",
                    "total": 2,
                    "mcq_count": 1,
                    "short_count": 1,
                    "mcq_options": 4,
                    "title_hint": "API Generated",
                },
            )

        self.assertTrue(generated["ok"], generated.get("errors"))
        self.assertTrue(generated["output_path"])
        output_path = Path(generated["output_path"])
        self.assertTrue(output_path.exists())
        self.assertEqual(output_path.parent.parent, (self.root / "Quizzes").resolve())

        pdf_materials = [m for m in generated["extracted_materials"] if m["path"].endswith(".pdf")]
        self.assertEqual(len(pdf_materials), 1)
        self.assertIn("needs_ocr", pdf_materials[0])
        if not pdf_materials[0]["needs_ocr"]:
            self.assertTrue(str(pdf_materials[0].get("content", "")).strip())

        bad_response = self.client.post(
            "/v1/generate/run",
            headers=self.headers,
            json={
                "sources": collect["sources"],
                "provider": "claude",
                "model": "stub-model",
                "total": 2,
                "mcq_count": 1,
                "short_count": 1,
                "mcq_options": 4,
                "title_hint": "Blocked",
                "output_subdir": "../outside",
            },
        )
        self.assertEqual(bad_response.status_code, 422, bad_response.text)
        self.assertIn("inside the managed Quizzes directory", bad_response.text)

    def test_collect_sources_can_return_extracted_materials_for_context_injection(self) -> None:
        docs = self.root / "docs"
        docs.mkdir(parents=True, exist_ok=True)
        txt_path = docs / "notes.txt"
        txt_path.write_text("hello from txt", encoding="utf-8")
        md_path = docs / "notes.md"
        md_path.write_text("# Heading\n\nhello from markdown", encoding="utf-8")

        docx_path = docs / "notes.docx"
        has_docx = False
        try:
            import docx

            doc = docx.Document()
            doc.add_paragraph("hello from docx")
            doc.save(str(docx_path))
            has_docx = True
        except Exception:
            has_docx = False

        pptx_path = docs / "slides.pptx"
        has_pptx = False
        try:
            from pptx import Presentation

            deck = Presentation()
            slide = deck.slides.add_slide(deck.slide_layouts[1])
            slide.shapes.title.text = "Deck Title"
            slide.placeholders[1].text = "Deck body"
            slide.notes_slide.notes_text_frame.text = "Speaker note"
            deck.save(str(pptx_path))
            has_pptx = True
        except Exception:
            has_pptx = False

        collect = self._post(
            "/v1/generate/collect-sources",
            {"paths": [str(docs)], "include_content": True},
        )
        self.assertIn("extracted_materials", collect)
        materials = collect["extracted_materials"]
        self.assertTrue(materials)

        by_name = {Path(item["path"]).name: item for item in materials}
        self.assertIn("notes.txt", by_name)
        self.assertIn("hello from txt", by_name["notes.txt"]["content"])
        self.assertIn("notes.md", by_name)
        self.assertIn("hello from markdown", by_name["notes.md"]["content"].lower())

        if has_docx:
            self.assertIn("notes.docx", by_name)
            self.assertIn("hello from docx", by_name["notes.docx"]["content"].lower())
        if has_pptx:
            self.assertIn("slides.pptx", by_name)
            pptx_text = str(by_name["slides.pptx"]["content"] or "")
            self.assertIn("Deck Title".lower(), pptx_text.lower())
            self.assertIn("Speaker note".lower(), pptx_text.lower())

    def test_generate_run_falls_back_when_requested_claude_model_unavailable(self) -> None:
        docs = self.root / "docs"
        docs.mkdir(parents=True, exist_ok=True)
        txt_path = docs / "notes.txt"
        txt_path.write_text("hello from txt", encoding="utf-8")

        collect = self._post(
            "/v1/generate/collect-sources",
            {"paths": [str(docs)]},
        )
        provider = _ModelAwareStubProvider(available_models=["claude-3-5-haiku-latest"])
        with patch("quiz_app.api.server._provider_client", return_value=provider):
            generated = self._post(
                "/v1/generate/run",
                {
                    "sources": collect["sources"],
                    "provider": "claude",
                    "model": "claude-3-7-sonnet-latest",
                    "total": 2,
                    "mcq_count": 1,
                    "short_count": 1,
                    "mcq_options": 4,
                    "title_hint": "Fallback test",
                },
            )

        self.assertTrue(generated["ok"], generated.get("errors"))
        self.assertEqual(provider.last_generation_model, "claude-3-5-haiku-latest")
        self.assertTrue(any("Falling back" in line for line in generated["warnings"]))

    def test_generate_run_prefers_preferred_model_key_for_claude_fallback(self) -> None:
        docs = self.root / "docs"
        docs.mkdir(parents=True, exist_ok=True)
        txt_path = docs / "notes.txt"
        txt_path.write_text("hello from txt", encoding="utf-8")

        self.client.put(
            "/v1/settings",
            headers=self.headers,
            json={"preferred_model_key": "claude:claude-3-opus-latest"},
        )

        collect = self._post(
            "/v1/generate/collect-sources",
            {"paths": [str(docs)]},
        )
        provider = _ModelAwareStubProvider(
            available_models=["claude-3-5-haiku-latest", "claude-3-opus-latest"]
        )
        with patch("quiz_app.api.server._provider_client", return_value=provider):
            generated = self._post(
                "/v1/generate/run",
                {
                    "sources": collect["sources"],
                    "provider": "claude",
                    "model": "claude-3-7-sonnet-latest",
                    "total": 2,
                    "mcq_count": 1,
                    "short_count": 1,
                    "mcq_options": 4,
                    "title_hint": "Preferred fallback test",
                },
            )

        self.assertTrue(generated["ok"], generated.get("errors"))
        self.assertEqual(provider.last_generation_model, "claude-3-opus-latest")
        self.assertTrue(any("Falling back" in line for line in generated["warnings"]))

    def test_generate_run_normalizes_mismatched_total(self) -> None:
        class _CaptureProvider:
            provider_name = "stub"

            def __init__(self) -> None:
                self.calls: list[dict[str, object]] = []

            def generate_quiz(
                self,
                materials_text: str,
                title_hint: str,
                instructions_hint: str,
                total_questions: int,
                mcq_count: int,
                short_count: int,
                mcq_options: int,
                model: str | None = None,
            ) -> Quiz:
                self.calls.append(
                    {
                        "materials_text": materials_text,
                        "title_hint": title_hint,
                        "instructions_hint": instructions_hint,
                        "total_questions": total_questions,
                        "mcq_count": mcq_count,
                        "short_count": short_count,
                        "mcq_options": mcq_options,
                        "model": model,
                    }
                )
                return Quiz(
                    title=title_hint or "Generated Quiz",
                    instructions=instructions_hint or "Answer all questions.",
                    questions=[
                        MCQQuestion(
                            id="q1",
                            prompt="Pick A",
                            points=1,
                            options=["A", "B", "C", "D"],
                            answer="A",
                        ),
                        MCQQuestion(
                            id="q2",
                            prompt="Pick B",
                            points=1,
                            options=["A", "B", "C", "D"],
                            answer="B",
                        ),
                    ],
                )

        docs = self.root / "docs"
        docs.mkdir(parents=True, exist_ok=True)
        txt_path = docs / "notes.txt"
        txt_path.write_text("hello from txt", encoding="utf-8")

        collect = self._post(
            "/v1/generate/collect-sources",
            {"paths": [str(docs)]},
        )

        provider = _CaptureProvider()
        with patch("quiz_app.api.server._provider_client", return_value=provider):
            generated = self._post(
                "/v1/generate/run",
                {
                    "sources": collect["sources"],
                    "provider": "claude",
                    "model": "stub-model",
                    "total": 99,
                    "mcq_count": 2,
                    "short_count": 0,
                    "mcq_options": 4,
                    "title_hint": "Normalized Total",
                },
            )

        self.assertTrue(generated["ok"], generated.get("errors"))
        self.assertTrue(generated["output_path"])
        self.assertTrue(any("normalized total" in line.lower() for line in generated["warnings"]))
        self.assertEqual(len(provider.calls), 1)
        call = provider.calls[0]
        self.assertEqual(call["total_questions"], 2)
        self.assertEqual(call["mcq_count"], 2)
        self.assertEqual(call["short_count"], 0)


if __name__ == "__main__":
    unittest.main()
